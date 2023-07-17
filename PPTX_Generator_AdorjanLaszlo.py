import collections.abc
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import json
from tkinter import Tk
from tkinter.filedialog import askopenfilename

# Create the Tkinter root window
root = Tk()
root.withdraw()  

# Prompt the user to choose the input .json file
file_path = askopenfilename()

def Generate_Picture(picname):
    x_values2 = [1, 10, 100, 1000, 10000, 100000, 1000000, 10000000]
    y_values2 = [320, 300, 260, 200, 151, 120, 95, 80]
    plt.figure(figsize=(10, 6))
    plt.plot(x_values2, y_values2, 'b')
    plt.plot(x_values2, y_values2, 'ro')
    plt.xscale('log')
    plt.xlim([1, 10000000])
    plt.ylim([0, 350])
    plt.grid(True, which='both')
    plt.xlabel('Life (cycles)')
    plt.ylabel('Stress (MPa)')
    plt.title('S-N Curve for brittle aluminium with a UTS of 320 MPa')
    plt.savefig(picname)
    plt.close()

def Generate_plot(source_data, output_pic, slide):
    # Collecting the x-y values from the .dat file
            x_values = []
            y_values = []
            with open(source_data) as file:
              for line in file:
                values = line.strip().split(';')
                x = float(values[0])
                y = float(values[1])
                x_values.append(x)
                y_values.append(y)

            # Plotting the x-y values and saving the plot as a .png
            plt.plot(x_values, y_values, 'b')
            plt.xlabel(slide['configuration']['x-label'])
            plt.ylabel(slide['configuration']['y-label'])
            plt.savefig(output_pic)
            plt.close()

# Reading the configuration file here

########################### Prompting the user here to select a folder
with open(file_path) as file:
    dict = json.load(file)

# Chosing the presentation from input and start iteration over the slides
pptx = dict['presentation']
prs = Presentation()
for slide in pptx:

    # Iterating through the keys of an individual slide
    for key in slide:
     if key == 'type':
        # The layout and content of the slide is dependent on the 'type' key. So, 5 different types are distinguished in the following.
        # Other config keys are adjusted after the 'type' key is found

        if slide[key].lower() == 'title':
            Slide_Layout = prs.slide_layouts[0]
            added_slide = prs.slides.add_slide(Slide_Layout)
            title = added_slide.shapes.title
            title.text = slide['title']           
            subtitle = added_slide.placeholders[1]
            subtitle.text = slide['content']

        if slide[key].lower() == 'text':
            Slide_Layout = prs.slide_layouts[5]
            added_slide = prs.slides.add_slide(Slide_Layout)
            title = added_slide.shapes.title
            title.text = slide['title']           
            txBox = added_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(7), Inches(10))
            tf = txBox.text_frame		
            tf.text = slide['content'] 

        if slide[key].lower() == 'list':   
            Slide_Layout = prs.slide_layouts[1]
            added_slide = prs.slides.add_slide(Slide_Layout)
            title = added_slide.shapes.title
            title.text = slide['title']           
            text_frame = added_slide.placeholders[1].text_frame

            # Iterating over the list elements and setting up the bullet points
            for new_line in slide['content']:
              new_paragraph = text_frame.add_paragraph()
              new_paragraph.text = new_line['text']
              new_paragraph.level = new_line['level']

        if slide[key].lower() == 'picture':
            Slide_Layout = prs.slide_layouts[5]
            added_slide = prs.slides.add_slide(Slide_Layout)
            title = added_slide.shapes.title
            title.text = slide['title']  
            
            # Generation of picture with a user-defined method 
            # The data points are not given, I have decided to generate the picture by reading the values from the sample pptx
            # When we actually have the picture,  it needs to be added to the slide directly in line #99
            Generate_Picture('SN_aluminium.png')
            added_slide.shapes.add_picture('SN_aluminium.png', Inches(0), Inches(1.2)) # added_slide.shapes.add_picture(slide['content'])   # Alternatively if input filepath is set properly
        
        if slide[key].lower() == 'plot':
            # Generation of plot with a user-defined method
            Generate_plot(slide['content'], 'plot.png', slide)

            # Importing the .png  plot into the slide
            Slide_Layout = prs.slide_layouts[5]
            added_slide = prs.slides.add_slide(Slide_Layout)
            title = added_slide.shapes.title
            title.text = slide['title']
            added_slide.shapes.add_picture('plot.png', Inches(1.5), Inches(1.5))

prs.save(r'Generated_presentation_AdorjanLaszlo.pptx')